"""Parse PPT files and UOR Excel files into structured data."""

import re
import io
import openpyxl
from pptx import Presentation


def _extract_shape_texts(shape, texts: list) -> None:
    """Recursively extract text from a shape, including grouped sub-shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Recurse into group shapes (the most common cause of missed content)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            _extract_shape_texts(sub_shape, texts)
        return

    # Text frames
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = " ".join(r.text.strip() for r in para.runs if r.text.strip())
            if line:
                texts.append(line)

    # Tables
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                texts.append(" | ".join(cells))


def extract_slides_text(pptx_bytes: bytes) -> dict[int, str]:
    """Return {slide_number: text} from a PPTX file.

    Recursively extracts text from grouped shapes, which is essential for
    slides built using Google Slides or PowerPoint shape groups — the most
    common reason slide content was silently missing.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = {}
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            _extract_shape_texts(shape, texts)
        if texts:
            slides[i] = "\n".join(texts)
        else:
            slides[i] = ""
    return slides


def get_slide_excerpt(slides: dict[int, str], start: int, end: int) -> str:
    """Get concatenated slide text for a slide range (inclusive)."""
    parts = []
    for num in range(start, end + 1):
        text = slides.get(num, "")
        if text:
            parts.append(f"[Slide {num}]\n{text}")
    return "\n\n".join(parts)


def parse_slide_range(range_str: str) -> tuple[int, int]:
    """Parse '4–13' or '4-13' or '4' into (start, end)."""
    if not range_str:
        return (0, 0)
    range_str = str(range_str).strip().replace("–", "-").replace("—", "-")
    if "-" in range_str:
        parts = range_str.split("-")
        try:
            return (int(parts[0].strip()), int(parts[-1].strip()))
        except ValueError:
            return (0, 0)
    try:
        n = int(range_str)
        return (n, n)
    except ValueError:
        return (0, 0)


def parse_uor_excel(xlsx_bytes: bytes) -> list[dict]:
    """
    Parse a CIBOP UOR Excel file into a list of UOR dicts.
    Each dict: {uor_id, title, objective, competency, sub_competencies,
                ears, slide_nos_source, extra}
    Sub-competencies with no UOR ID are attached to the previous UOR as extras.
    """
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
    ws = wb.active

    # Find header row
    header_row = None
    header_map = {}
    for i, row in enumerate(ws.iter_rows(values_only=True), start=1):
        if row and any(str(c or "").strip().lower().startswith("uor") for c in row):
            header_row = i
            for j, cell in enumerate(row):
                if cell:
                    header_map[str(cell).strip().lower()[:30]] = j
            break

    if header_row is None:
        return []

    def col(row_vals, key_fragment):
        """Get value from row by partial header key match."""
        for k, idx in header_map.items():
            if key_fragment.lower() in k:
                v = row_vals[idx] if idx < len(row_vals) else None
                return str(v).strip() if v else ""
        return ""

    uors = []
    current_uor = None

    for row in ws.iter_rows(min_row=header_row + 1, values_only=True):
        if not any(c for c in row):
            continue

        uor_id = col(row, "uor id") or col(row, "uorid")
        title = col(row, "uor title") or col(row, "title")
        objective = col(row, "objective")
        competency = col(row, "competency")
        sub_comp = col(row, "sub-comp") or col(row, "subcomp")
        ears = col(row, "ear")
        slide_src = col(row, "source ppt") or col(row, "slide nos")

        if uor_id and uor_id.upper() != "UOR ID":
            # Parse sub-competencies from the cell (newline or bullet separated)
            sub_list = _parse_sub_competencies(sub_comp)
            current_uor = {
                "uor_id": uor_id,
                "title": title,
                "objective": objective,
                "competency": competency,
                "sub_competencies": sub_list,
                "ears": ears,
                "slide_nos_source": slide_src,
                "extra_scs": []
            }
            uors.append(current_uor)
        elif current_uor is not None and title:
            # Sub-row with no UOR ID — attach as extra SC to previous UOR
            extra_sc = {
                "title": title,
                "objective": objective,
                "sub_competency": sub_comp,
                "ears": ears,
            }
            current_uor.setdefault("extra_scs", []).append(extra_sc)

    return uors


def _parse_sub_competencies(text: str) -> list[str]:
    """Split sub-competency text into individual SC strings."""
    if not text:
        return []
    # Split on newlines or SC patterns like SC1.1, SC2.1 etc
    parts = re.split(r'\n|(?=SC\d+\.\d+)', text)
    result = []
    for p in parts:
        p = p.strip()
        if p:
            result.append(p)
    return result


def extract_pdf_text(pdf_bytes: bytes) -> str:
    """Extract text from PDF bytes using markitdown or fallback."""
    try:
        import tempfile, os
        from markitdown import MarkItDown
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(pdf_bytes)
            tmp_path = f.name
        md = MarkItDown()
        result = md.convert(tmp_path)
        os.unlink(tmp_path)
        return result.text_content if result else ""
    except Exception:
        # Fallback: try pypdf
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""
