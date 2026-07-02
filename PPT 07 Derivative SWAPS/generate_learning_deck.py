#!/usr/bin/env python3
"""Generate SWP Learning Deck PPTX"""
import subprocess, sys
subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'], check=True)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

OUTPUT_PATH = '/Users/sandeeppr/CIBOP/PPT 07 Derivative SWAPS/SWP_Learning_Deck_v2.pptx'

# Colors
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
TEAL    = RGBColor(0x00, 0x5F, 0x5F)
GOLD    = RGBColor(0x8B, 0x69, 0x14)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF2, 0xF6, 0xFB)
LTEAL   = RGBColor(0xE0, 0xF2, 0xF2)
LGOLD   = RGBColor(0xFD, 0xF5, 0xE0)
DARKGRAY= RGBColor(0x33, 0x33, 0x33)

W = Inches(10)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # blank

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_name='Calibri', font_size=18, bold=False, italic=False,
                color=DARKGRAY, align=PP_ALIGN.LEFT, wrap=True, v_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_textbox(slide, bullets, left, top, width, height,
                       font_name='Calibri', font_size=16, color=DARKGRAY,
                       title=None, title_color=TEAL):
    from pptx.util import Pt as Pt2
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = font_name
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for bullet in bullets:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f'  •  {bullet}'
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        first = False
    return txBox

def slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ──────────────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────────────

def make_hook_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, NAVY)

    # Gold accent bar top
    add_rect(slide, 0, 0, W, Inches(0.12), fill_color=GOLD)

    # Module tag
    add_textbox(slide, 'CIBOP | CAPITAL MARKETS & INVESTMENT BANKING OPERATIONS PROGRAMME',
                Inches(0.5), Inches(0.2), Inches(9), Inches(0.4),
                font_size=9, color=GOLD, bold=True, align=PP_ALIGN.LEFT)

    # Main title
    add_textbox(slide, 'DERIVATIVES: SWAPS',
                Inches(0.5), Inches(1.2), Inches(9), Inches(1.4),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, 'Module 7 — From Plain Vanilla IRS to Credit Default Swaps',
                Inches(0.5), Inches(2.5), Inches(8.5), Inches(0.7),
                font_size=18, color=TEAL, align=PP_ALIGN.LEFT)

    # Hook statement
    add_rect(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(2.2), fill_color=TEAL)
    add_textbox(slide, '"Every day, trillions of dollars in interest rate risk, currency exposure, and credit risk change hands through one instrument that most people have never heard of — the swap. By the end of this module, you will be the person who has."',
                Inches(0.7), Inches(3.5), Inches(8.6), Inches(2.0),
                font_size=15, italic=True, color=WHITE, align=PP_ALIGN.LEFT, wrap=True)

    # Gold bottom bar
    add_rect(slide, 0, Inches(7.3), W, Inches(0.2), fill_color=GOLD)
    add_textbox(slide, 'SWP.1 — SWP.6  |  6 Units of Readiness  |  Interest Rate | Currency | Credit Default Swaps',
                Inches(0.5), Inches(6.9), Inches(9), Inches(0.35),
                font_size=9, color=GOLD, align=PP_ALIGN.LEFT)

# ──────────────────────────────────────────────────────
# Per-UOR slide maker helpers
# ──────────────────────────────────────────────────────

def uor_title_bar(slide, uor_id, uor_title, bg_color=NAVY):
    """Dark top bar with UOR ID and title"""
    add_rect(slide, 0, 0, W, Inches(1.0), fill_color=bg_color)
    add_textbox(slide, uor_id, Inches(0.35), Inches(0.08), Inches(1.2), Inches(0.45),
                font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_textbox(slide, uor_title, Inches(0.35), Inches(0.45), Inches(9.3), Inches(0.5),
                font_size=19, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def slide_subtitle(slide, subtitle, top=Inches(1.1), color=TEAL):
    add_textbox(slide, subtitle, Inches(0.4), top, Inches(9.2), Inches(0.45),
                font_size=13, bold=True, color=color, align=PP_ALIGN.LEFT)

def make_quick_recap(prs, uor_id, uor_title, recap_points, prior_uor_title):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT)
    uor_title_bar(slide, uor_id, uor_title)
    slide_subtitle(slide, f'Quick Recap — Before we start {uor_id}', top=Inches(1.05))

    add_textbox(slide, f'Building on: {prior_uor_title}',
                Inches(0.4), Inches(1.55), Inches(9.2), Inches(0.35),
                font_size=11, italic=True, color=DARKGRAY)

    add_bullet_textbox(slide, recap_points,
                       Inches(0.5), Inches(2.0), Inches(9.0), Inches(4.5),
                       font_size=16, color=DARKGRAY)

    # Teal left accent bar
    add_rect(slide, Inches(0.3), Inches(2.0), Inches(0.06), Inches(4.5), fill_color=TEAL)

def make_scenario(prs, uor_id, uor_title, scenario_title, scenario_body):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LGOLD)
    uor_title_bar(slide, uor_id, uor_title, bg_color=TEAL)
    slide_subtitle(slide, f'Scenario — {scenario_title}', color=NAVY)

    # Scenario card
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(9.0), Inches(4.9), fill_color=WHITE)
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(0.1), Inches(4.9), fill_color=GOLD)

    add_textbox(slide, scenario_body,
                Inches(0.75), Inches(1.85), Inches(8.6), Inches(4.6),
                font_size=15, color=DARKGRAY, wrap=True)

def make_key_concept(prs, uor_id, uor_title, concept_title, left_bullets, right_note):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT)
    uor_title_bar(slide, uor_id, uor_title)
    slide_subtitle(slide, f'Key Concept — {concept_title}')

    # Left column — bullets
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(5.5), Inches(5.4), fill_color=WHITE)
    add_bullet_textbox(slide, left_bullets,
                       Inches(0.55), Inches(1.75), Inches(5.25), Inches(5.2),
                       font_size=15, color=DARKGRAY)

    # Right column — note/chart box
    add_rect(slide, Inches(6.1), Inches(1.65), Inches(3.5), Inches(5.4), fill_color=LTEAL)
    add_textbox(slide, right_note,
                Inches(6.2), Inches(1.75), Inches(3.3), Inches(5.2),
                font_size=13, color=NAVY, wrap=True)

def make_going_deeper(prs, uor_id, uor_title, deeper_title, deeper_bullets):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT)
    uor_title_bar(slide, uor_id, uor_title, bg_color=TEAL)
    slide_subtitle(slide, f'Going Deeper — {deeper_title}', color=NAVY)

    add_rect(slide, Inches(0.4), Inches(1.65), Inches(9.2), Inches(5.4), fill_color=WHITE)
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(0.08), Inches(5.4), fill_color=GOLD)

    add_bullet_textbox(slide, deeper_bullets,
                       Inches(0.6), Inches(1.75), Inches(8.9), Inches(5.2),
                       font_size=15, color=DARKGRAY)

def make_think_discuss(prs, uor_id, uor_title, questions):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LGOLD)
    uor_title_bar(slide, uor_id, uor_title, bg_color=GOLD)
    add_textbox(slide, 'Think & Discuss',
                Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.45),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    for i, q in enumerate(questions):
        top = Inches(1.7 + i * 1.7)
        add_rect(slide, Inches(0.4), top, Inches(9.2), Inches(1.5), fill_color=WHITE)
        add_rect(slide, Inches(0.4), top, Inches(0.1), Inches(1.5), fill_color=NAVY)
        add_textbox(slide, f'Q{i+1}  {q}',
                    Inches(0.65), Inches(top/914400 * 14400 / 12700 + 0.08), Inches(8.8), Inches(1.35),
                    font_size=14, color=DARKGRAY, wrap=True)

def make_interview_ready(prs, uor_id, uor_title, qa_pairs, next_uor_label=None):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT)
    uor_title_bar(slide, uor_id, uor_title, bg_color=NAVY)
    add_textbox(slide, 'Interview Ready — Q&A',
                Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.45),
                font_size=14, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    row_h = Inches(1.45)
    for i, (q, a) in enumerate(qa_pairs):
        top = Inches(1.6) + i * row_h
        # Q box
        add_rect(slide, Inches(0.4), top, Inches(3.8), row_h - Inches(0.1), fill_color=NAVY)
        add_textbox(slide, f'Q: {q}',
                    Inches(0.5), top + Inches(0.1), Inches(3.6), row_h - Inches(0.25),
                    font_size=12, bold=True, color=WHITE, wrap=True)
        # A box
        add_rect(slide, Inches(4.4), top, Inches(5.2), row_h - Inches(0.1), fill_color=LTEAL)
        add_textbox(slide, f'A: {a}',
                    Inches(4.5), top + Inches(0.1), Inches(5.0), row_h - Inches(0.25),
                    font_size=12, color=NAVY, wrap=True)

    if next_uor_label:
        add_rect(slide, Inches(0.4), Inches(6.8), Inches(9.2), Inches(0.45), fill_color=GOLD)
        add_textbox(slide, f'Next Unit: {next_uor_label}',
                    Inches(0.5), Inches(6.82), Inches(9.0), Inches(0.4),
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────
# SLIDE DATA PER UOR
# ──────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ── COVER ──
    make_hook_cover(prs)

    # ══════════════════════════════════════════
    # SWP.1 — What is a Swap?
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.1', 'What is a Swap?'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'Derivatives derive value from an underlying asset — stocks, rates, currencies, commodities',
            'Types of derivatives: Forwards, Futures, Options, Swaps',
            'OTC derivatives are bilateral contracts — not exchange-traded; terms are customisable',
            'Hedging: reducing existing risk | Speculation: taking on risk for profit',
        ],
        prior_uor_title='Basics of Derivatives (DV.1–DV.3)')

    make_scenario(prs, uid, utitle,
        scenario_title='The Rate Mismatch Problem',
        scenario_body='Infra Corp has borrowed Rs 500 crore at a floating rate (MIBOR + 1%). Their revenue is predictable — toll collections from a highway. Their CFO cannot sleep: "If MIBOR rises 2%, our interest cost jumps Rs 10 crore a year. Our revenue does not move."\n\nTech Dynamics has the opposite problem. They issued fixed-rate bonds at 8%, but they are in a growth phase — they expect RBI to cut rates significantly. They are locked in at 8% while new borrowers will pay 6%.\n\nQuestion: Can Infra Corp and Tech Dynamics help each other?\n\nAnswer: Yes. Through a swap — an agreement to exchange their cash flow obligations. Each gets what the other has. No loan transfer. No refinancing. Just a contractual exchange of payment streams. This is the founding logic of the $700 trillion global swap market.')

    make_key_concept(prs, uid, utitle,
        concept_title='The Swap as a Cash Flow Exchange',
        left_bullets=[
            'A swap is an OTC agreement between two parties to exchange a series of cash flows based on different terms over an agreed period',
            'Two counterparties: Fixed Payer (pays set rate) and Floating Payer (pays MIBOR-based rate)',
            'Notional principal = reference amount ONLY — never physically exchanged',
            'Governed by ISDA Master Agreement — bilateral, customisable, not exchange-traded',
            'Key uses: Hedging rate risk | Speculation on rate direction | Comparative advantage arbitrage',
            'Economically equivalent to a series of forward rate agreements, one per settlement date',
        ],
        right_note='[DIAGRAM PLACEHOLDER]\n\nParty A (Fixed Payer)\n→ Pays 7% fixed →\n← Receives MIBOR ←\nParty B (Floating Payer)\n\nNotional: Rs 100 Crore\nTenor: 3 Years\nSettlement: Semi-annual\nNet: Only difference changes hands\n\n[See Source Slide 3-5 — Swap Structure Diagram]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='Comparative Advantage — Why Swaps Create Real Value',
        deeper_bullets=[
            'Comparative advantage: each party borrows where they are relatively cheapest, then swaps to desired rate type',
            "Example: Infra Co borrows fixed at 8%; Tech Co at 10%. Infra Co borrows floating at MIBOR+1%; Tech Co at MIBOR+2%",
            'Infra Co has comparative advantage in fixed markets; Tech Co in floating — difference = 2% spread advantage',
            'Each borrows where cheapest, then swaps: both end up with preferred rate at LOWER all-in cost than direct borrowing',
            'Total saving = 2% per annum on notional — divided between parties via swap pricing',
            'This is not a zero-sum game — it genuinely creates value by arbitraging different credit perceptions in different markets',
            'Classic example from Salomon Brothers structuring the first currency swap for IBM and World Bank (1981)',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'A friend argues: "Swaps are just gambling — you\'re betting on which way rates move." How would you respond? Is there a distinction between using swaps for hedging vs speculation, and does that distinction matter morally?',
            'The notional principal is never exchanged in a plain vanilla IRS, yet a swap is described as a "Rs 1,000 crore IRS." Why does the notional matter if it never moves?',
            'If comparative advantage in swaps creates real value for both parties, who is on the other side bearing the cost? Is someone always losing?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('What is a swap and how does it differ from a futures contract?',
             'A swap is an OTC agreement to exchange cash flows over a period — customisable, bilateral, ISDA-governed. Futures are standardised, exchange-traded, marked daily. A swap is like a tailored suit; a future is off-the-rack.'),
            ('Why is the notional principal never exchanged in an IRS?',
             'The notional is purely a reference amount to calculate interest. Only the net interest difference flows — this is why a Rs 100 crore swap involves only Rs 0.5–1 crore of actual cash movement per period, not Rs 100 crore.'),
            ('What is the comparative advantage argument for entering a swap?',
             'Each party borrows where they have relatively lower cost, then swaps. If Infra Co is better in fixed markets and Tech Co in floating, total borrowing cost is lower than if each borrowed directly in their preferred market.'),
        ],
        next_uor_label='SWP.2 — Interest Rate Swaps (IRS)')

    # ══════════════════════════════════════════
    # SWP.2 — Interest Rate Swaps (IRS)
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.2', 'Interest Rate Swaps (IRS)'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'A swap = OTC agreement to exchange cash flows based on different terms',
            'Two parties: fixed payer and floating payer; notional never exchanged',
            'Uses: hedging interest rate risk, speculation, comparative advantage',
            'ISDA Master Agreement governs; swap is a series of forward rate agreements',
        ],
        prior_uor_title='SWP.1 — What is a Swap?')

    make_scenario(prs, uid, utitle,
        scenario_title='The Manufacturing CFO\'s Rate Nightmare',
        scenario_body='National Manufacturing Ltd borrowed Rs 200 crore from SBI at MIBOR + 1.5%. The loan runs for 5 years. Current MIBOR: 6.5%. All-in cost: 8%.\n\nThe CFO builds their capital budget around Rs 16 crore of annual interest. But the RBI is in a hiking cycle. If MIBOR hits 9%, interest jumps to Rs 21 crore. That Rs 5 crore swing is the difference between profit and loss.\n\nSolution: Enter a plain vanilla IRS. Pay fixed 8% to a bank; receive MIBOR from the bank. The bank loan still charges MIBOR + 1.5% — but the IRS receipt of MIBOR cancels that floating component. Net cost: 8% fixed + 1.5% spread = 9.5% total — locked. Certain. Predictable.\n\nThe CFO sleeps. The RBI hikes. The swap pays off — MIBOR rises to 9.5%, the IRS generates Rs 3 crore of net receipt, partially offsetting the spread. The hedge works exactly as designed.')

    make_key_concept(prs, uid, utitle,
        concept_title='Plain Vanilla IRS Mechanics',
        left_bullets=[
            'Fixed payer pays agreed fixed rate; floating payer pays MIBOR (resets each period)',
            'MIBOR = Mumbai Interbank Offer Rate — Indian equivalent of LIBOR',
            'Notional: Rs 100 crore | Fixed: 7% p.a. | Semi-annual settlement',
            'Period calculation: Fixed = 7% × 100Cr × 0.5 = Rs 3.5 Cr per period',
            'Floating = MIBOR × 100Cr × 0.5 (varies each period)',
            'Net settlement: only the DIFFERENCE is paid — no gross exchange',
            'MIBOR 6%: Fixed payer pays Rs 0.5 Cr | MIBOR 8%: Fixed payer receives Rs 0.5 Cr',
        ],
        right_note='[CALCULATION TABLE]\nNotional: Rs 100 Cr\nFixed: 7% p.a., semi-annual\n\nPeriod 1 (MIBOR=6%):\nFixed: Rs 3.5 Cr\nFloating: Rs 3.0 Cr\nNet: Fixed pays Rs 0.5 Cr\n\nPeriod 2 (MIBOR=8%):\nFixed: Rs 3.5 Cr\nFloating: Rs 4.0 Cr\nNet: Floating pays Rs 0.5 Cr\n\nPeriod 3 (MIBOR=7%):\nNet: Zero\n\n[See Source Slide 8 — IRS Cash Flow Table]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='Who Benefits When Rates Move — The Rate Bet Inside Every IRS',
        deeper_bullets=[
            'Pay-fixed party benefits when rates RISE above the fixed rate — receive floating income grows, fixed cost stays same',
            'Pay-floating party benefits when rates FALL below the fixed rate — pay floating falls, fixed receipt stays same',
            'This makes IRS a natural rate speculation vehicle for macro traders and hedge funds',
            'RBI hiking cycle → macro fund enters pay-fixed OIS; if MIBOR rises to 8.5% vs fixed 7% → profit on each settlement',
            'Corporate hedger enters pay-fixed to REMOVE the speculation — they want certainty, not the rate bet',
            'OIS (Overnight Index Swap) is the most liquid IRS in India — based on MIBOR overnight, not term MIBOR',
            'India IRS market turnover: over Rs 2 lakh crore daily; most active OTC derivative segment in India (RBI data)',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'A corporate CFO says: "We entered a pay-fixed IRS to hedge our floating loan, but now rates have fallen and we\'re losing money on the swap." Has the hedge failed? How do you explain this to the board?',
            'MIBOR is the floating benchmark in India, but LIBOR was phased out globally in 2021. What challenges do you think arose when trillions of dollars of swap contracts referenced LIBOR as their floating rate?',
            'If a bank can simply pass its interest rate risk to a corporate through a swap, who ultimately absorbs that risk in the system?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('Walk me through the cash flows of a plain vanilla IRS.',
             'Party A pays 7% fixed on Rs 100 crore semi-annually = Rs 3.5 crore. Party B pays MIBOR on same notional. Net: only the difference flows. If MIBOR is 6%, A pays Rs 0.5 crore. If MIBOR is 8%, B pays Rs 0.5 crore. Notional never moves.'),
            ('Why do corporates use IRS to hedge floating-rate loans?',
             'A corporate with floating debt (MIBOR+1.5%) enters pay-fixed IRS at 7%. The MIBOR they receive from the swap offsets the MIBOR they pay on the loan. Net cost becomes fixed at 7% + 1.5% = 8.5% — regardless of where MIBOR goes. Certainty for budgeting.'),
            ('What is MIBOR and why is it the floating rate benchmark in Indian IRS?',
             'MIBOR = Mumbai Interbank Offer Rate — the rate Indian banks charge each other for short-term unsecured lending. It reflects RBI policy rate expectations and market liquidity. It resets daily or periodically, making it the natural floating benchmark for Indian OTC derivatives.'),
        ],
        next_uor_label='SWP.3 — Currency Swaps')

    # ══════════════════════════════════════════
    # SWP.3 — Currency Swaps
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.3', 'Currency Swaps'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'Plain vanilla IRS: fixed ↔ floating on same currency notional; only net interest flows',
            'MIBOR is the Indian floating rate benchmark; resets each period',
            'Pay-fixed party benefits when rates rise; pay-floating benefits when rates fall',
            'Corporates use IRS to convert floating debt to fixed — locking predictable costs',
        ],
        prior_uor_title='SWP.2 — Interest Rate Swaps (IRS)')

    make_scenario(prs, uid, utitle,
        scenario_title='Indian IT Company\'s Dollar Dilemma',
        scenario_body='TechSoft India Ltd needs $100 million to fund a US acquisition. Options:\n\n1. Borrow in India at 9% in INR, then convert — but INR loan is expensive and FX risk on repayment remains.\n\n2. Issue dollar bonds in the US at 4% (US investors trust investment-grade Indian IT) — but then service a $4M annual dollar coupon from INR revenues, bearing full FX risk.\n\n3. Do option 2, then enter a currency swap: receive $100M, pay INR equivalent (Rs 700 crore at 70 $/INR). Pay 7% INR interest to swap counterparty. Receive 4% USD interest from counterparty — exactly offsetting the dollar bond coupon.\n\nResult: TechSoft accesses $100M at effectively 7% INR cost — cheaper than domestic 9% — with NO currency mismatch. The acquisition is funded. The CFO\'s FX exposure is eliminated. The currency swap did the heavy lifting.')

    make_key_concept(prs, uid, utitle,
        concept_title='Currency Swap Structure — Three Phases',
        left_bullets=[
            'Phase 1 — Inception: Principal exchanged at current spot rate (Rs 700 Cr ↔ $100M at 70)',
            'Phase 2 — Periodic: Each party pays interest in the RECEIVED currency',
            '  Party A pays USD 4% on $100M = $4M annually to counterparty',
            '  Party B pays INR 7% on Rs 700 Cr = Rs 49 Cr annually to Party A',
            'Phase 3 — Maturity: Principals returned at the ORIGINAL rate of 70 — regardless of spot',
            'FX risk on principal: ELIMINATED (original rate locks the exchange permanently)',
            'Key difference from IRS: Principal IS physically exchanged — creating larger counterparty exposure',
        ],
        right_note='[DIAGRAM PLACEHOLDER]\n\nInception:\nParty A → Rs 700 Cr → Party B\nParty B → $100M → Party A\n\nPeriodic:\nParty A → $4M → Party B\nParty B → Rs 49 Cr → Party A\n\nMaturity:\nParty A → $100M → Party B\nParty B → Rs 700 Cr → Party A\n(at ORIGINAL rate of 70)\n\n[See Source Slide 15 — Currency Swap Flow]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='Central Bank Currency Swaps and Sovereign Applications',
        deeper_bullets=[
            'RBI uses bilateral currency swap lines with Bank of Japan (USD 75 billion SAARC swap) and others',
            'Purpose: shore up FX reserves during stress without selling existing reserves',
            'During COVID-19, RBI drew on these lines to provide dollar liquidity to Indian banks',
            'World Bank and IBM executed the first currency swap in 1981 — IBM had DM/CHF debt it wanted in USD',
            'Currency swap maturity risk: principal exchange at maturity creates settlement risk across two settlement systems',
            'ISDA/CSA collateral frameworks: because principal is exchanged, variation margin may not fully cover exposure',
            'Post-trade: currency swaps must be reported to CCIL trade repository; central clearing for standardised tenors is developing in India',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'A CFO says: "Our currency swap eliminates all our FX risk." Is this entirely accurate? What residual risks remain even after a well-structured currency swap?',
            'The original exchange rate is locked at maturity in a currency swap, regardless of where spot USD/INR is then. If the rupee depreciates sharply, who benefits and who bears the cost in a bilateral currency swap?',
            'Why might a central bank prefer a currency swap line with another central bank over simply buying foreign currency in the spot market?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('How does a currency swap differ from an interest rate swap?',
             'IRS: same currency, no principal exchange, only net interest flows. Currency swap: two currencies, principal IS exchanged at inception and maturity at the original rate. Currency swap hedges both interest rate AND FX risk simultaneously.'),
            ('Why is principal exchanged in a currency swap but not in an IRS?',
             'Because the two parties need the other\'s currency to service their obligations. An Indian company that issued USD bonds needs to pay dollars — so it physically receives dollars via the swap. The principal exchange is functional, not just notional.'),
            ('What risk does a currency swap NOT eliminate?',
             'Counterparty risk — because principal is physically exchanged, if the counterparty defaults after you have delivered your currency but before maturity, you may not recover the other currency. FX risk on the principal is eliminated; credit risk is amplified.'),
        ],
        next_uor_label='SWP.4 — Credit Default Swaps (CDS)')

    # ══════════════════════════════════════════
    # SWP.4 — Credit Default Swaps (CDS)
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.4', 'Credit Default Swaps (CDS)'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'Currency swap: exchanges principal in two currencies at inception and maturity at original rate',
            'Eliminates FX risk on principal; periodic interest paid in received currency',
            'Corporate application: access cheap foreign capital; swap to domestic currency obligations',
            'Larger counterparty risk than IRS because principal physically moves',
        ],
        prior_uor_title='SWP.3 — Currency Swaps')

    make_scenario(prs, uid, utitle,
        scenario_title='The Bond Holder\'s Insurance Question',
        scenario_body='HDFC Bank holds Rs 500 crore of Yes Bank bonds in its trading book. Year is 2019. Yes Bank is under stress — governance concerns, NPA disclosures, promoter pledging. HDFC\'s credit analysts are worried: "What if Yes Bank defaults?"\n\nOption 1: Sell the bonds. But selling Rs 500 crore of distressed bonds in a thin market will move prices badly. HDFC takes a mark-to-market loss AND signals to the market that they are exiting.\n\nOption 2: Buy CDS protection. Pay, say, 300 basis points per annum (Rs 15 crore annually) to a protection seller. If Yes Bank triggers a credit event — bankruptcy, payment failure, restructuring — the protection seller pays HDFC the face value of the bonds.\n\nHDFC keeps the bonds (no market impact), pays the insurance premium (Rs 15 Cr/year), and sleeps knowing the downside is capped. This is CDS in its legitimate, intended use: transferring credit risk without disturbing the underlying market.')

    make_key_concept(prs, uid, utitle,
        concept_title='CDS Structure — Insurance on a Bond',
        left_bullets=[
            'Protection Buyer: pays CDS spread (e.g. 150 bps p.a.) quarterly — like an insurance premium',
            'Protection Seller: collects premium; pays face value of bond if credit event occurs',
            'Reference Entity: the company/sovereign whose credit risk is being transferred',
            'Credit Events (ISDA-defined): Bankruptcy, Failure to Pay, Debt Restructuring, Repudiation/Moratorium',
            'Settlement: Physical (deliver defaulted bond, receive face value) or Cash (pay recovery-adjusted difference)',
            'CDS Spread = market\'s real-time price of default probability and severity',
            'Higher spread = higher default risk priced in (AA-rated: 20–50 bps; distressed: 500+ bps)',
        ],
        right_note='[DIAGRAM PLACEHOLDER]\n\nProtection BUYER\n→ pays 150 bps p.a. quarterly →\n← receives face value on default ←\nProtection SELLER\n\nReference Entity:\n[Company / Sovereign]\n\nCredit Event Examples:\n• Bankruptcy filing\n• Missed coupon payment\n• Forced debt restructuring\n• Sovereign moratorium\n\nCDS Spread timeline:\nNormal: 80 bps\nStress: 400 bps\nPre-default: 1500 bps\n\n[See Source Slide 20 — CDS Structure]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='CDS and the 2008 Crisis — From Hedge to Weapon',
        deeper_bullets=[
            'AIG Financial Products sold $500 billion+ in CDS protection on CDOs backed by US subprime mortgages',
            'Implicit assumption: US housing prices would never fall broadly — correlation was underestimated',
            'When housing fell, ALL CDS triggered simultaneously — correlated defaults, not independent events',
            '"Naked" CDS: buyers held no underlying bonds — pure speculation on default; amplified systemic pressure on Lehman, Bear Stearns',
            'AIG could not post collateral as CDO values fell; US government $182B bailout to prevent AIG collapse from triggering cascade',
            'Post-2008 reforms: Dodd-Frank/EMIR mandated central clearing for standardised CDS; ISDA Determinations Committee formalised',
            'EU banned naked sovereign CDS in 2012; ongoing debate about whether CDS can destabilise sovereign bond markets',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'AIG sold CDS protection without sufficient capital reserves. What risk management failure does this represent? How should a CDS seller manage their exposure?',
            '"Naked CDS is like buying fire insurance on your neighbour\'s house — you have an incentive for it to burn." Do you agree? Should naked CDS be banned globally?',
            'The CDS spread on a sovereign bond is often described as more informative than credit ratings. Why might market-implied spreads be more accurate than agency ratings in predicting default?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('What is a Credit Default Swap and how does it work?',
             'CDS = insurance on a bond. Protection buyer pays regular spread (premium). Protection seller pays face value if a credit event (bankruptcy, payment failure, restructuring) occurs. The spread reflects the market\'s view of default probability.'),
            ('What triggered the 2008 financial crisis in the context of CDS?',
             'AIG sold massive CDS protection on CDOs (bundled subprime mortgages) without adequate capital. When US housing fell broadly, all CDOs defaulted simultaneously. AIG could not pay — $182B government bailout was required. Naked CDS amplified panic around Lehman and Bear Stearns.'),
            ('What is the difference between a naked CDS and a hedging CDS?',
             'Hedging CDS: buyer owns the underlying bond and uses CDS to cap downside — legitimate risk transfer. Naked CDS: buyer has no underlying exposure — pure speculation on default. Same instrument, entirely different economic function and systemic risk implications.'),
        ],
        next_uor_label='SWP.5 — Swap Pricing & Valuation Basics')

    # ══════════════════════════════════════════
    # SWP.5 — Swap Pricing & Valuation Basics
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.5', 'Swap Pricing & Valuation Basics'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'CDS = insurance on a bond; protection buyer pays spread; seller pays on credit event',
            'CDS spread = market\'s real-time price of default probability',
            '2008 crisis: AIG sold unhedged CDS; correlated defaults; systemic failure',
            'Post-2008: central clearing, ISDA determinations committee, EU ban on naked sovereign CDS',
        ],
        prior_uor_title='SWP.4 — Credit Default Swaps (CDS)')

    make_scenario(prs, uid, utitle,
        scenario_title='The Swap Book Manager\'s Daily Reality',
        scenario_body='Priya runs the interest rate derivatives desk at a large Indian bank. She has 150 swap positions on her book. Every morning, she gets a report: each swap\'s mark-to-market value, the net MTM across the book, and the collateral calls that result.\n\nToday, RBI surprised the market with a 50 bps rate cut. Priya\'s pay-fixed swaps (swaps where she pays fixed, receives floating) are now showing negative MTM — she agreed to receive MIBOR, which just fell 50 bps, but she still pays the same fixed rate.\n\nHer team needs to answer: How much negative value has been generated? How much variation margin must she post to CCPs today? Which positions should she hedge with offsetting trades?\n\nAll of these questions require swap valuation — the ability to calculate the present value of remaining cash flows, adjusted for the new rate environment. This is swap pricing in practice.')

    make_key_concept(prs, uid, utitle,
        concept_title='Zero Value at Inception — and How That Changes',
        left_bullets=[
            'At inception: swap has ZERO value — swap rate set so PV(fixed) = PV(floating)',
            'Swap rate derived from yield curve — the fixed rate that makes the trade "fair"',
            'In India: OIS curve (overnight index swap) used to derive discount factors',
            'After execution: value changes every time rates move, curve shifts, or time passes',
            'MTM (Mark-to-Market) = PV(remaining floating cash flows) - PV(remaining fixed cash flows)',
            'Rates rise → positive MTM for pay-fixed party (fixed payment below market)',
            'Rates fall → negative MTM for pay-fixed party (fixed payment above market)',
        ],
        right_note='[VALUATION DIAGRAM]\n\nAt inception:\nPV(fixed leg) = PV(floating leg)\nSwap value = 0\n\nAfter rates rise 1%:\nPV(floating) increases\nPay-fixed MTM = positive\n\nAfter rates fall 1%:\nPV(floating) decreases\nPay-fixed MTM = negative\n\nMTM drivers:\n1. Rate level change\n2. Yield curve shape\n3. Time passage\n4. Counterparty credit\n\n[See Source Slide 24 — Swap Valuation]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='CVA, DVA and the Three Ways to Exit a Swap',
        deeper_bullets=[
            'CVA (Credit Valuation Adjustment): discounts swap MTM for probability that YOUR COUNTERPARTY defaults before maturity',
            'DVA (Debt Valuation Adjustment): adjustment for your OWN default probability — creates perverse P&L when your credit worsens',
            'Banks must calculate CVA/DVA on all OTC derivative books — significant capital charge under Basel III',
            'Exit Option 1: Offsetting swap — enter reverse position (pay floating, receive fixed); neutralises P&L without closing original',
            'Exit Option 2: Early termination — mutually agree to cancel; pay/receive current MTM as cash settlement',
            'Exit Option 3: Novation — assign your swap position to a third party (requires original counterparty consent)',
            'Central clearing (CCIL): variation margin posted daily = swap value essentially marked daily, reducing credit risk',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'A new analyst asks: "If the swap is worth zero at inception, how does the bank make money structuring and selling swaps to corporate clients?" How would you explain the bank\'s economics on swap transactions?',
            'When a pay-fixed IRS shows negative MTM for the corporate (because rates fell), the corporate\'s CFO says "our swap is losing money — we should terminate it." Should the CFO terminate? What are the arguments for and against?',
            'CVA means that the value of a derivative depends partly on the creditworthiness of the counterparty. Does this mean a derivative with a technically insolvent counterparty has zero value? How do banks manage this?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('Why does a swap have zero value at inception?',
             'The swap rate is chosen to make PV(fixed leg) = PV(floating leg), using the current yield curve. Both parties are exchanging equal present values — so neither has an advantage at inception. Zero value = fair price.'),
            ('How does a swap\'s MTM value change after inception?',
             'When market rates rise, the fixed payer\'s position gains positive MTM — they pay below-market fixed, receive above-market floating. When rates fall, MTM turns negative. Four drivers: rate level, yield curve shape, time decay, counterparty credit quality (CVA).'),
            ('What is CVA and why does it matter for swap valuation?',
             'CVA = Credit Valuation Adjustment. It reduces the MTM of a derivative by the probability that the counterparty defaults before maturity. A swap with positive MTM owed by a shaky counterparty is worth less than face value. Basel III requires banks to hold capital against CVA risk.'),
        ],
        next_uor_label='SWP.6 — Swap Applications & Risk Management')

    # ══════════════════════════════════════════
    # SWP.6 — Swap Applications & Risk Management
    # ══════════════════════════════════════════
    uid, utitle = 'SWP.6', 'Swap Applications & Risk Management'

    make_quick_recap(prs, uid, utitle,
        recap_points=[
            'Swap rate makes PV(fixed) = PV(floating) at inception → zero value',
            'MTM changes as rates move: pay-fixed gains when rates rise, loses when rates fall',
            'CVA adjusts for counterparty credit risk; DVA for own credit',
            'Three exit options: offsetting swap, early termination (pay MTM), novation',
        ],
        prior_uor_title='SWP.5 — Swap Pricing & Valuation Basics')

    make_scenario(prs, uid, utitle,
        scenario_title='The Treasury Head\'s Strategic Swap Decision',
        scenario_body='Axis Bank\'s treasury head faces a classic ALM (Asset-Liability Management) challenge:\n\nAssets: Rs 10,000 crore of 5-year fixed-rate home loans at 8.5% — long-term, fixed rate\nLiabilities: Rs 8,000 crore of 6-month FDs at MIBOR + 0.5% — short-term, floating rate\n\nIf RBI hikes rates: deposit costs rise immediately; loan income stays fixed at 8.5% → NIM compressed.\n\nSolution: Enter a receive-fixed IRS. Receive 8% fixed, pay MIBOR floating.\nNow: Rising rates hurt deposits (pay more MIBOR on FDs) but the swap generates floating payments TO Axis (MIBOR income from IRS). The hedge offsets the mismatch.\n\nThis is textbook Asset-Liability Management using swaps — managing the interest rate sensitivity of the balance sheet without changing the underlying loan or deposit terms.')

    make_key_concept(prs, uid, utitle,
        concept_title='Swap Applications Across User Types',
        left_bullets=[
            'Banks (ALM): Receive-fixed IRS hedges fixed-rate loan / floating-rate deposit mismatch',
            'Corporate (liability management): Pay-fixed converts floating debt to fixed — lock in cost',
            'Corporate (liability management): Pay-floating converts fixed bonds to floating — benefit from rate cuts',
            'Macro fund (speculation): Pay-fixed OIS to bet on RBI rate hikes; receive-fixed to bet on cuts',
            'Currency swap: Indian corporate accesses cheaper USD capital, swaps to INR obligations',
            'CDS: Bank buys protection to reduce credit exposure on corporate loan book without selling loans',
            'Cross-currency IRS: Combines currency swap and IRS in one instrument — swaps both currency and rate type',
        ],
        right_note='[APPLICATION MATRIX]\n\nUser | Instrument | Purpose\nBank Treasury | Receive-fixed IRS | ALM / NIM protection\nCorporate CFO | Pay-fixed IRS | Hedge floating debt\nCorporate CFO | Pay-floating IRS | Liability restructuring\nMacro Fund | Pay-fixed OIS | Speculate on rate rise\nIT Company | Currency swap | USD bonds → INR obligation\nCred. Bank | CDS (buy) | Reduce loan book exposure\n\n[See Source Slide 30 — Applications Table]')

    make_going_deeper(prs, uid, utitle,
        deeper_title='Central Clearing, Hedge Accounting and the Full Risk Spectrum',
        deeper_bullets=[
            'CCIL (Clearing Corporation of India) — clears standardised IRS and OIS centrally; eliminates bilateral counterparty risk',
            'Central clearing mechanics: CCP interposes as buyer/seller; initial margin + daily variation margin required',
            'Hedge accounting (Ind AS 109/IFRS 9): Cash flow hedge → effective portion to OCI (not P&L); reduces P&L volatility',
            'Fair value hedge → gains/losses on swap and hedged item both to P&L simultaneously — offset each other',
            'Swap risks: Market risk (rate/FX changes MTM) | Counterparty risk | Liquidity risk (OTC, harder to exit) | Operational risk',
            'Counterparty risk larger in long-dated currency swaps — principal physically exchanged over long tenor',
            'Regulatory capital: Basel III charges capital for counterparty credit risk (CCR) and CVA on uncleared swaps',
        ])

    make_think_discuss(prs, uid, utitle,
        questions=[
            'A bank\'s ALM head says: "We don\'t need to use swaps — we can just match our loan and deposit maturities by changing our product mix." Is this a viable alternative to swap hedging? What are the trade-offs?',
            'Under IFRS 9 hedge accounting, a corporate must document the hedge relationship formally before entering the trade. Why do you think regulators require this documentation? What might go wrong without it?',
            'Swaps are described as both a risk management tool and a source of systemic risk (as seen in 2008). Can both be true? What structural features determine whether a swap reduces or amplifies systemic risk?',
        ])

    make_interview_ready(prs, uid, utitle,
        qa_pairs=[
            ('How do banks use IRS to manage interest rate risk on their balance sheet?',
             'Banks have fixed-rate loan assets funded by floating-rate deposits — a structural mismatch. They enter receive-fixed IRS: receive fixed (offsetting fixed loan income), pay floating (matching deposit costs). If rates rise, IRS floating payments increase, offsetting higher deposit costs.'),
            ('What is the comparative advantage argument for currency swaps?',
             'Each party borrows where relatively cheaper and swaps. If Indian co is better in INR markets and US co in USD, each borrows in home market (where they have the edge) and swaps. Total borrowing cost for both is lower than if each borrowed directly in the desired currency.'),
            ('What are the main risks in a swap position?',
             'Market risk (MTM moves with rates/FX), counterparty risk (default before maturity), liquidity risk (OTC = harder to exit than exchange-traded), and operational risk (settlement, documentation). Counterparty risk is largest in long-dated currency swaps where principal is exchanged.'),
        ],
        next_uor_label=None)  # Module complete

    # ── Final module complete slide ──
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, NAVY)
    add_rect(slide, 0, 0, W, Inches(0.12), fill_color=GOLD)
    add_textbox(slide, 'MODULE COMPLETE',
                Inches(0.5), Inches(1.5), Inches(9), Inches(0.8),
                font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 'DERIVATIVES: SWAPS — SWP.1 to SWP.6',
                Inches(0.5), Inches(2.4), Inches(9), Inches(0.6),
                font_size=20, color=GOLD, align=PP_ALIGN.CENTER)
    coverage = [
        'SWP.1  What is a Swap?  —  OTC cash flow exchange, ISDA, fixed/floating parties',
        'SWP.2  Interest Rate Swaps  —  Plain vanilla IRS, MIBOR, net settlement mechanics',
        'SWP.3  Currency Swaps  —  Principal exchange, FX risk elimination, cross-border funding',
        'SWP.4  Credit Default Swaps  —  CDS structure, credit events, 2008 crisis lessons',
        'SWP.5  Swap Pricing & Valuation  —  Zero inception value, MTM, CVA, exit routes',
        'SWP.6  Swap Applications  —  ALM, liability management, speculation, central clearing',
    ]
    for i, line in enumerate(coverage):
        add_textbox(slide, line,
                    Inches(1.5), Inches(3.2 + i * 0.55), Inches(7.5), Inches(0.5),
                    font_size=12, color=LIGHT, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, Inches(7.3), W, Inches(0.2), fill_color=GOLD)

    prs.save(OUTPUT_PATH)
    print(f'Saved: {OUTPUT_PATH}')
    import os
    print(f'File size: {os.path.getsize(OUTPUT_PATH):,} bytes')
    print(f'Slides: {len(prs.slides)}')

if __name__ == '__main__':
    build_deck()
    print('Learning Deck DONE.')
