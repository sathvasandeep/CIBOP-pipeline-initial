#!/usr/bin/env python3
"""Generate STO & RDM Learning Deck — CIBOP Module 8"""
import subprocess, sys
subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'], check=True)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUTPUT_PATH = '/Users/sandeeppr/CIBOP/PPT 08 Securities/STO_Learning_Deck_v2.pptx'

# ─────────────────────────────────────────────
# COLORS
# ─────────────────────────────────────────────
NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
TEAL  = RGBColor(0x00, 0x5F, 0x5F)
GOLD  = RGBColor(0x8B, 0x69, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF8, 0xFF)
LGREY = RGBColor(0xEE, 0xEE, 0xEE)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────
# UTILITIES
# ─────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, bullets, x, y, w, h, size=16, color=DARK, title=None, title_color=NAVY):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(size + 2)
        r.font.bold = True
        r.font.color.rgb = title_color
    for i, bullet in enumerate(bullets):
        if title and i == 0:
            p = tf.add_paragraph()
        elif not title and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 1
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f'• {bullet}'
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return txBox

def add_slide_header(slide, uor_id, uor_title, slide_label, bg_color=NAVY):
    add_rect(slide, 0, 0, 13.33, 1.1, bg_color)
    add_text_box(slide, f'{uor_id} — {uor_title}', 0.3, 0.05, 9.5, 0.5,
                 size=14, bold=True, color=WHITE)
    add_text_box(slide, slide_label, 0.3, 0.55, 9.5, 0.45,
                 size=11, color=GOLD, bold=False)
    add_text_box(slide, 'CIBOP | Module 8 | STO & RDM', 10.0, 0.05, 3.2, 0.4,
                 size=9, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
# SLIDE TYPES
# ─────────────────────────────────────────────

def make_hook_cover(prs):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 5.8, 13.33, 1.7, TEAL)
    add_text_box(slide, 'CIBOP', 0.5, 0.3, 4.0, 0.6, size=13, color=GOLD, bold=True)
    add_text_box(slide, 'MODULE 8', 0.5, 0.9, 5.0, 0.5, size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
    add_text_box(slide,
        'Securities Trading\nOrganisation & RDM',
        0.5, 1.6, 12.0, 2.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide,
        'Thousands of trades execute every second. Behind every one sits an invisible machine — '
        'the Securities Trading Organisation. This module unpacks that machine: who does what, '
        'how trades settle, and why the data beneath it all is the difference between a clean book and a catastrophic fail.',
        0.5, 3.9, 10.5, 1.8, size=14, color=RGBColor(0xCC, 0xDD, 0xFF), italic=True)
    add_text_box(slide, 'STO.1 – STO.6', 0.5, 6.0, 5.0, 0.5, size=13, color=WHITE, bold=True)
    add_text_box(slide, 'Front • Middle • Back Office • Reference Data • Corporate Actions',
                 0.5, 6.5, 10.0, 0.6, size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

def make_quick_recap(prs, uor_id, uor_title, bullets):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    add_slide_header(slide, uor_id, uor_title, 'Quick Recap', bg_color=NAVY)
    add_rect(slide, 0.3, 1.25, 12.73, 0.04, TEAL)
    add_bullet_box(slide, bullets, 0.5, 1.4, 12.0, 5.8,
                   size=17, color=DARK,
                   title='Before we begin — key ideas from prior learning:',
                   title_color=TEAL)

def make_scenario(prs, uor_id, uor_title, scenario_title, scenario_text, discussion_q):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    add_slide_header(slide, uor_id, uor_title, 'Scenario / Case Study', bg_color=TEAL)
    add_rect(slide, 0.3, 1.2, 12.73, 0.04, GOLD)
    add_text_box(slide, scenario_title, 0.4, 1.35, 12.5, 0.55,
                 size=16, bold=True, color=NAVY)
    add_text_box(slide, scenario_text, 0.4, 1.95, 12.3, 3.4,
                 size=14, color=DARK, wrap=True)
    add_rect(slide, 0.3, 5.5, 12.73, 1.7, NAVY)
    add_text_box(slide, f'💡  Think about it: {discussion_q}',
                 0.5, 5.6, 12.2, 1.5, size=13, color=WHITE, italic=True)

def make_key_concept(prs, uor_id, uor_title, concept_title, concept_points, placeholder_note=None):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    add_slide_header(slide, uor_id, uor_title, 'Key Concept', bg_color=NAVY)
    add_rect(slide, 0.3, 1.2, 12.73, 0.04, GOLD)
    add_text_box(slide, concept_title, 0.4, 1.3, 12.5, 0.6, size=18, bold=True, color=NAVY)
    add_bullet_box(slide, concept_points, 0.5, 2.0, 12.0, 4.5, size=15, color=DARK)
    if placeholder_note:
        add_rect(slide, 0.4, 6.1, 12.5, 1.1, LGREY)
        add_text_box(slide, placeholder_note, 0.5, 6.15, 12.2, 1.0,
                     size=10, color=RGBColor(0x55, 0x55, 0x55), italic=True)

def make_going_deeper(prs, uor_id, uor_title, deeper_title, deeper_points):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK)
    add_slide_header(slide, uor_id, uor_title, 'Going Deeper', bg_color=TEAL)
    add_rect(slide, 0.3, 1.2, 12.73, 0.04, GOLD)
    add_text_box(slide, deeper_title, 0.4, 1.3, 12.5, 0.6,
                 size=18, bold=True, color=GOLD)
    add_bullet_box(slide, deeper_points, 0.5, 2.0, 12.0, 5.0,
                   size=14, color=WHITE)

def make_think_discuss(prs, uor_id, uor_title, questions):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    add_slide_header(slide, uor_id, uor_title, 'Think & Discuss', bg_color=GOLD)
    add_rect(slide, 0.3, 1.2, 12.73, 0.04, NAVY)
    add_text_box(slide, 'Discussion Questions', 0.4, 1.3, 12.0, 0.55,
                 size=16, bold=True, color=NAVY)
    for i, q in enumerate(questions):
        y = 1.95 + i * 1.5
        add_rect(slide, 0.4, y, 12.5, 1.3, RGBColor(0xE8, 0xEF, 0xF8))
        add_text_box(slide, f'Q{i+1}.  {q}', 0.6, y + 0.1, 12.1, 1.1,
                     size=14, color=DARK, wrap=True)

def make_interview_ready(prs, uor_id, uor_title, qa_pairs):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_slide_header(slide, uor_id, uor_title, 'Interview Ready', bg_color=TEAL)
    add_rect(slide, 0.3, 1.2, 12.73, 0.04, GOLD)
    add_text_box(slide, '🎯  Questions You Must Be Able to Answer', 0.4, 1.3, 12.5, 0.55,
                 size=16, bold=True, color=GOLD)
    for i, (question, answer) in enumerate(qa_pairs):
        y = 2.0 + i * 1.6
        add_rect(slide, 0.4, y, 12.5, 1.45, RGBColor(0x0A, 0x1E, 0x3A))
        add_text_box(slide, f'Q: {question}', 0.55, y + 0.05, 12.1, 0.5,
                     size=12, bold=True, color=GOLD)
        add_text_box(slide, f'A: {answer}', 0.55, y + 0.55, 12.1, 0.85,
                     size=11, color=WHITE, wrap=True)

# ─────────────────────────────────────────────
# MODULE CONTENT
# ─────────────────────────────────────────────

make_hook_cover(prs)

# ── STO.1 ─────────────────────────────────────
ID, TITLE = 'STO.1', 'Introduction to the Securities Trading Organisation'

make_quick_recap(prs, ID, TITLE, [
    'Capital markets involve the buying and selling of securities — equities, bonds, derivatives.',
    'A trade is not complete at execution — it must be confirmed, cleared, and settled.',
    'Multiple parties are involved: the bank, counterparty, exchange, clearinghouse, and custodian.',
    'Operational failures — settlement fails, reconciliation breaks — carry financial and regulatory costs.',
])

make_scenario(prs, ID, TITLE,
    'Day in the Life: Trade 1234 Meets the STO',
    'At 10:03 AM, a fund manager instructs CIBOP Bank to buy 500,000 shares of HDFC Bank at market. '
    'The order enters the OMS, passes pre-trade compliance checks, and hits the NSE in 200 milliseconds. '
    'Fill confirmed. Now the back office has until T+1 to confirm the trade with the counterparty, instruct the custodian, '
    'and ensure the securities arrive in the client\'s account. Four systems, twelve people, and one golden rule: '
    'no settlement failure.',
    'If any one of the three STO divisions — front, middle, back — makes an error, who catches it, and what is the cost?'
)

make_key_concept(prs, ID, TITLE,
    'The Three Divisions of the STO',
    [
        'FRONT OFFICE: Revenue-generating — traders, sales, research, structuring. Executes trades and manages client relationships.',
        'MIDDLE OFFICE: Governance layer — risk management, P&L attribution, independent price verification, compliance.',
        'BACK OFFICE: Settlement engine — confirms trades, instructs custodians, reconciles positions, processes corporate actions.',
        'CHINESE WALL: Mandatory information barrier between front office and back/middle office — prevents insider trading and front-running.',
        'STP RATE: Straight-Through Processing — the % of trades that settle without manual intervention. Target: >95%.',
    ],
    '[DIAGRAM: See Source Slide — STO Three-Division org chart with role examples in each division]'
)

make_going_deeper(prs, ID, TITLE,
    'STP Rate, Settlement Metrics & Systems',
    [
        'Settlement fail rate: <1% is the industry benchmark; CSDR/SEBI impose financial penalties above threshold.',
        'OMS (Order Management System) → EMS (Execution) → Risk System → Settlement System: each feeds the next.',
        'A break in the system chain causes manual intervention — raising cost-per-trade and operational risk.',
        'Cost per trade = total operational cost / trade volume — the metric that drives automation investment.',
        'Reconciliation breaks aged >3 days are escalated to management and flagged in audit/compliance reports.',
        'The Chinese Wall is enforced via physical separation, IT access controls, and email monitoring policies.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'An STP rate of 92% sounds impressive — but at a bank processing 200,000 trades per day, how many manual interventions does that represent? What are the cost and risk implications?',
    'The Chinese Wall is a regulatory requirement, but does it ever create operational problems? How do banks manage legitimate communication needs across the wall?',
    'If you had to improve one STO performance metric — STP rate, settlement fail rate, or reconciliation breaks — which would you choose and why?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What are the three main divisions of a Securities Trading Organisation?',
     'Front Office (revenue/execution), Middle Office (risk/governance/P&L), Back Office (operations/settlement/custody/reconciliation).'),
    ('What is Straight-Through Processing (STP) and why does it matter?',
     'STP is automated trade processing from execution to settlement with zero manual intervention. A high STP rate reduces cost, operational risk, and settlement fails.'),
    ('What is the Chinese Wall and why is it legally required?',
     'An information barrier preventing front office traders from accessing client or settlement data in the back/middle office — required to prevent insider trading and front-running.'),
])

# ── STO.2 ─────────────────────────────────────
ID, TITLE = 'STO.2', 'Front Office Operations'

make_quick_recap(prs, ID, TITLE, [
    'The front office generates revenue through trading, sales, research, and structuring.',
    'Orders originate from clients, flow through sales to traders, and are executed on market.',
    'The OMS manages orders pre- and post-execution; the EMS routes orders to exchanges.',
    'Best execution is a regulatory obligation under MiFID II and SEBI rules.',
])

make_scenario(prs, ID, TITLE,
    'The Wrong-Side Trade',
    'At 2:15 PM, an equity trader at CIBOP Bank executes what she believes is a client sell order for 1 million '
    'Infosys shares. She enters the trade as a BUY instead. The error is not caught until the middle office P&L '
    'attribution review at end of day — by which point Infosys has moved Rs 12 per share. The bank is now long '
    '1 million shares it should not hold, with a Rs 1.2 crore market risk exposure. The error log is filed, '
    'the risk desk is alerted, and the client is never told — but the compliance team certainly is.',
    'What controls should have caught this before end of day? And who bears the cost of the error?'
)

make_key_concept(prs, ID, TITLE,
    'Order Flow & Best Execution',
    [
        'Order Flow: Client → Sales (order received, logged) → Trader (execution instruction with algo choice) → OMS → EMS → Exchange → Fill confirmation → Back to OMS.',
        'Pre-trade Compliance Checks: OMS automatically checks restricted list, position limits, account eligibility before releasing any order to market.',
        'VWAP / TWAP / IS Algorithms: execution strategies that balance market impact against speed — selected based on order size, urgency, and market conditions.',
        'Best Execution (MiFID II / SEBI): firms must achieve the best possible result for clients — price, speed, likelihood of execution, size, cost. All documented in a Best Execution Policy.',
        'Agency vs Proprietary: agency = no bank risk (earn commission); proprietary = bank\'s own capital at risk (post-Volcker: heavily restricted).',
    ],
    '[DIAGRAM: See Source Slide — Order flow diagram from client instruction through OMS/EMS to execution]'
)

make_going_deeper(prs, ID, TITLE,
    'Front Office P&L and Error Management',
    [
        'Daily P&L = Realised (closed positions) + Unrealised MTM (open positions) − Commissions − Financing costs.',
        'Traders do NOT mark their own books — middle office IPV uses independent pricing sources to verify every position.',
        'P&L Attribution breaks down daily results: market move delta P&L, carry/theta, new trades, unexplained residual.',
        'Unexplained P&L triggers investigation before books close — common causes: missing trades, data error, pricing discrepancy.',
        'Error types: wrong security (ISIN error), wrong side (most expensive), wrong quantity, duplicate trade — all require immediate escalation.',
        'Error P&L impact threshold: typically £10k–£100k before mandatory senior management escalation.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'A trader argues that best execution is always about getting the best price. The regulations say otherwise — what other factors count, and why?',
    'The Volcker Rule restricted proprietary trading in banks. Was this the right response to the 2008 crisis, or did it remove a valuable market-making function?',
    'If the OMS pre-trade compliance check flags an order as breaching a position limit, should the trader be able to override it? What governance should surround any override?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What is the difference between agency trading and proprietary trading?',
     'Agency: bank executes on a client\'s behalf, earns a commission, bears no market risk. Proprietary: bank trades its own capital, bears full market risk. Volcker Rule restricts prop trading at deposit-taking banks.'),
    ('What does best execution mean under MiFID II?',
     'Firms must take all sufficient steps to obtain the best possible result for clients, considering price, speed, likelihood of execution, size, and cost. Documented in a Best Execution Policy.'),
    ('What is P&L attribution and why is it done independently?',
     'Decomposing daily P&L into market move, carry, new trades, and unexplained residual. Done independently by the middle office using separate pricing sources to prevent traders inflating their own reported profits.'),
])

# ── STO.3 ─────────────────────────────────────
ID, TITLE = 'STO.3', 'Middle Office Operations'

make_quick_recap(prs, ID, TITLE, [
    'The middle office provides the governance layer between front office and back office.',
    'It does not execute or settle trades — it validates, verifies, monitors, and escalates.',
    'Every trading desk operates within risk limits approved by the risk committee.',
    'IPV (Independent Price Verification) ensures P&L cannot be inflated by trader self-marking.',
])

make_scenario(prs, ID, TITLE,
    'The IPV Dispute',
    'CIBOP\'s structured credit trader marks a portfolio of illiquid CLO tranches at 98.5 — his model-based '
    'fair value. The middle office IPV team checks three independent pricing vendors: Bloomberg prices them at '
    '94.2 on average. The difference: Rs 42 crore on a Rs 1,000 crore book. The middle office applies an '
    'IPV reserve of Rs 42 crore against the desk\'s reported P&L. The trader protests. The risk committee '
    'reviews. The reserve stands. The trader\'s bonus is materially affected.',
    'Without IPV, how could a trader systematically inflate their P&L over time? What would eventually expose them?'
)

make_key_concept(prs, ID, TITLE,
    'Risk Limits, IPV, and P&L Attribution',
    [
        'Risk Limits: position limits (size), VaR limits (max expected daily loss at 95%/99% confidence), DV01 limits (fixed income rate sensitivity), counterparty and concentration limits.',
        'VaR (Value at Risk): statistical measure of the maximum expected loss over a given period at a given confidence level — e.g., "1-day 99% VaR = Rs 5 crore" means only a 1% chance of losing more than Rs 5 crore in one day.',
        'IPV (Independent Price Verification): middle office values all positions using external sources — Bloomberg, Refinitiv, broker quotes. Material differences → IPV reserve deducted from P&L.',
        'P&L Attribution: Delta P&L (price move × position) + Carry/Theta + New Trades P&L. Unexplained residual > threshold triggers investigation.',
        'Collateral (ISDA/CSA): daily variation margin calls issued and received based on net MTM exposure of OTC derivative positions.',
    ],
    '[TABLE: See Source Slide — Risk limit framework table showing limit types, measurement, and breach escalation process]'
)

make_going_deeper(prs, ID, TITLE,
    'Trade Confirmation & Exception Management',
    [
        'Affirmation: both sides agree on economic terms post-execution (instrument, direction, quantity, price, settlement date).',
        'Formal Confirmation: legal documents exchanged via DTCC TIW / Markit MatchIT for OTC trades; must match within T+1.',
        'Unmatched confirms aged >T+1 = compliance breach and settlement risk — must be escalated immediately.',
        'Exception report runs daily: unconfirmed trades, limit breaches, IPV differences, collateral disputes — each time-stamped and assigned to an owner.',
        'Initial Margin (IM): separate from variation margin — required under post-2008 BCBS/IOSCO rules for non-centrally cleared OTC derivatives.',
        'CVA (Credit Valuation Adjustment): reduces the value of your OTC derivative receivables to reflect counterparty default probability.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'A trader whose P&L is reduced by an IPV reserve argues the reserve is "conservative" and will "reverse out" next quarter. Under what circumstances might that be true, and when is it a red flag?',
    'VaR limits are widely used but widely criticised — they worked poorly in the 2008 crisis. What are VaR\'s known limitations, and what complementary risk measures are used alongside it?',
    'The middle office is independent of the front office — but they must work together daily. How do you maintain independence without creating a culture of adversarial checking?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What is Independent Price Verification (IPV) and why is it important?',
     'Using external, independent pricing sources to value the firm\'s trading positions, separate from the trader\'s own marks. Prevents P&L inflation and ensures reported profits are based on fair market values.'),
    ('What types of risk limits does the middle office enforce?',
     'Position limits (size), VaR limits (statistical max loss), DV01 limits (rate sensitivity), counterparty limits (max exposure per firm), and concentration limits (max sector exposure).'),
    ('What is variation margin in the context of OTC derivatives?',
     'Daily cash or eligible securities posted between ISDA/CSA counterparties to cover the net MTM exposure of OTC derivative positions. If your swap is in-the-money by Rs 10 crore, your counterparty posts Rs 10 crore to you.'),
])

# ── STO.4 ─────────────────────────────────────
ID, TITLE = 'STO.4', 'Back Office Operations'

make_quick_recap(prs, ID, TITLE, [
    'Settlement is when a trade becomes real: securities and cash change hands.',
    'India operates T+1 equity settlement — one of the fastest globally since January 2023.',
    'DVP (Delivery versus Payment) ensures simultaneous exchange — eliminating principal risk.',
    'Custodians hold and service securities; CSDs (NSDL, CDSL) are the book of record.',
])

make_scenario(prs, ID, TITLE,
    'The SSI That Wasn\'t Updated',
    'CIBOP Bank sells Rs 500 crore of Indian government bonds to a Japanese asset manager. '
    'The back office generates a settlement instruction using the SSI on file — which was last updated '
    'two years ago, before the Japanese firm changed its custodian. The instruction routes to the old '
    'custodian, who rejects it. It is now T+1 at 3:45 PM. The correct custodian\'s cut-off is 4:00 PM. '
    'The operations team has 15 minutes to locate the new SSI, re-instruct, and re-confirm. '
    'They miss the cut-off by 4 minutes. The trade fails to settle. Penalty accrues. The relationship manager gets the call.',
    'What would you put in place to prevent this SSI failure? And at what point in the process should the error have been caught?'
)

make_key_concept(prs, ID, TITLE,
    'Settlement: DVP, T+1, and Fail Management',
    [
        'DVP (Delivery versus Payment): settlement is atomic — securities and cash move simultaneously. No securities without cash; no cash without securities. Eliminates Herstatt Risk.',
        'Settlement Cycle: India T+1 (since Jan 2023), US T+1 (since May 2024), Europe T+2 for most equities. Shorter cycle = lower counterparty risk but faster operations required.',
        'Fail Causes: incorrect SSI, insufficient securities in custodian account, cash shortfall on value date, system/communication failure, unresolved trade dispute.',
        'Fails Management: Day 1 fail → chase counterparty; Day 4+ → escalation, potential buy-in notice. Buy-in: exchange purchases securities at failing party\'s cost.',
        'Three-Way Reconciliation: Internal records ↔ Custodian statement ↔ CSD records. Daily. Breaks must be resolved — aged breaks (>3 days) escalated.',
    ],
    '[DIAGRAM: See Source Slide — DVP settlement infrastructure diagram showing buyer, seller, custodians, and CSD]'
)

make_going_deeper(prs, ID, TITLE,
    'Custody, Asset Servicing & Reconciliation in Depth',
    [
        'Global Custodian: holds securities in legally segregated accounts, separate from firm\'s own assets (client money rules).',
        'Sub-Custodians: local market specialists appointed by the global custodian for markets without direct presence.',
        'Asset Servicing: dividend collection, coupon payments, corporate action processing, proxy voting, tax reclaim.',
        'Nostro Reconciliation: cash account reconciliation — comparing internal nostro ledger against the correspondent bank\'s statement for every currency.',
        'CSDR (Central Securities Depositories Regulation): EU regulation mandating settlement fails penalties and mandatory buy-ins for persistently failing trades.',
        'SEBI\'s framework: fines for settlement fails; auction mechanism for securities delivery failures in Indian equity markets.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'India moved to T+1 settlement — but critics argued it created operational challenges for foreign investors in different time zones. What are the pros and cons of shorter settlement cycles?',
    'A buy-in is described as an "extreme measure" — yet regulators in Europe made it mandatory under CSDR for persistent fails. Is making buy-ins mandatory a good regulatory response?',
    'Reconciliation is performed daily across three sets of records. What happens to the reconciliation process when a trade is disputed? How does the back office manage an open dispute while still keeping its books clean?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What is DVP and why is it important?',
     'Delivery versus Payment — a settlement mechanism where securities and cash are transferred simultaneously, so neither party delivers without receiving. Eliminates Herstatt Risk (principal loss from one-sided delivery).'),
    ('What are the main causes of settlement failures?',
     'Incorrect SSI routing, insufficient securities in the custodian account, cash shortfall on value date, missed system cut-off times, and unresolved trade disputes.'),
    ('What is three-way reconciliation?',
     'Comparing the firm\'s internal position records against the custodian\'s statement and the CSD\'s records daily. Differences (breaks) must be investigated and resolved — aged breaks are escalated.'),
])

# ── STO.5 ─────────────────────────────────────
ID, TITLE = 'STO.5', 'Reference Data Management (RDM)'

make_quick_recap(prs, ID, TITLE, [
    'Reference data is the static, foundational information about instruments, counterparties, and settlement routes.',
    'It underpins every stage of the trade lifecycle — from execution to settlement to reporting.',
    'Key types: Security Master (ISIN, asset class), Counterparty Data (LEI, BIC), SSI, Corporate Actions.',
    'Poor data quality cascades: one bad ISIN → wrong booking → wrong settlement → reconciliation break.',
])

make_scenario(prs, ID, TITLE,
    'The Duplicate ISIN Problem',
    'A corporate bond issued by Tata Steel has just been reopened — a new tranche with the same terms '
    'but a new ISIN. The reference data team at CIBOP is notified by Bloomberg but the update is delayed '
    'by two days. During those two days, three trades are booked using the old ISIN. Settlement instructions '
    'go to the wrong CSD participant. Two trades fail. One confirmation is rejected by the counterparty. '
    'Total cost: Rs 4.2 lakh in fail penalties plus half a day of operations time to manually fix each trade. '
    'The root cause: no automated feed from the CSD to the security master.',
    'What controls would you implement to prevent a reference data gap between a new bond issuance and the time it appears in the trading system?'
)

make_key_concept(prs, ID, TITLE,
    'The Golden Copy and Reference Data Lifecycle',
    [
        'Golden Copy: a single, authoritative, validated reference data record for each entity — consumed by all downstream systems. Eliminates version drift.',
        'Security Master: stores ISIN/CUSIP, asset class, currency, exchange, tick size, maturity — for every instrument the firm trades.',
        'Counterparty Data: LEI (Legal Entity Identifier — global unique ID), SWIFT BIC, legal name, KYC/AML status, credit limit.',
        'SSI (Standard Settlement Instructions): per-counterparty, per-currency, per-asset-class settlement routing — custodian, account, depository. Maintained via DTCC ALERT/Omgeo.',
        'Reference Data Lifecycle: Create → Validate (vs external sources) → Four-Eyes Approve → Distribute → Amend (controlled change) → Retire (not delete — preserve audit trail).',
    ],
    '[TABLE: See Source Slide — Reference data hierarchy table showing source authority (Bloomberg vs CSD vs SWIFT) by data type]'
)

make_going_deeper(prs, ID, TITLE,
    'SSI Deep Dive and Data Governance',
    [
        'SSI failure is the #1 cause of preventable settlement fails — stale SSI routes instructions to the wrong custodian or account.',
        'DTCC ALERT (Institutional): industry-managed SSI database where counterparties publish and maintain their official settlement instructions.',
        'Omgeo CTM / ALERT: matching platform used for SSI retrieval and trade confirmation in the institutional space.',
        'Data Governance: policies defining who can create/amend reference data, what the approval hierarchy is, and how conflicts between sources are resolved.',
        'GLEIF (Global Legal Entity Identifier Foundation): maintains the global LEI database — required for regulatory reporting under MiFID II, EMIR, CFTC.',
        'Corporate Actions as Reference Data: dividend dates, ex-dates, merger ratios, rights issue terms — all stored in the reference data system and consumed by asset servicing teams.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'The golden copy concept sounds ideal — but what happens when two authoritative sources (Bloomberg and the CSD) disagree on an ISIN or a corporate action date? How should the hierarchy be determined?',
    'Reference data teams are often understaffed compared to front office and even back office teams. How would you make the business case for increased investment in RDM?',
    'DTCC ALERT allows counterparties to self-publish their SSI data. What are the risks of a self-publication model, and how should a bank verify data it receives from a counterparty\'s own SSI submission?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What is reference data and how does it differ from market data?',
     'Reference data is static, foundational information about instruments and parties — ISINs, LEIs, SSIs, corporate action terms. Market data is real-time pricing and transactional data. Reference data errors cascade downstream; market data errors affect pricing.'),
    ('What is the golden copy concept in reference data management?',
     'A single, authoritative, validated record for each data entity — consumed by all downstream systems. Eliminates the version drift that occurs when each system maintains its own copy.'),
    ('Why are Standard Settlement Instructions (SSIs) so critical, and how are they managed?',
     'SSIs define exactly where securities and cash should be delivered for each counterparty. Stale or incorrect SSIs are the #1 cause of preventable settlement fails. Managed via industry platforms like DTCC ALERT/Omgeo, where counterparties publish their official instructions.'),
])

# ── STO.6 ─────────────────────────────────────
ID, TITLE = 'STO.6', 'Corporate Actions Processing'

make_quick_recap(prs, ID, TITLE, [
    'Corporate actions are company-declared events that affect the value or structure of securities.',
    'Mandatory events happen automatically; voluntary events require shareholder elections.',
    'Key dates: Declaration → Ex-Date → Record Date → Payment Date.',
    'Back office must identify eligible holders, notify clients, process elections, and update all position records.',
])

make_scenario(prs, ID, TITLE,
    'The Rights Issue That Went Wrong',
    'CIBOP Custody Bank holds 8.5 million shares of a mid-cap pharmaceutical company on behalf of 340 '
    'institutional clients. The company announces a rights issue: 1 new share for every 4 held, at a '
    '25% discount. Back office sends notifications — but three clients never receive theirs due to '
    'an outdated email address in the entity database. Two clients miss the election deadline and '
    'their rights lapse. Combined entitlement value lost: Rs 2.8 crore. Both clients file complaints. '
    'The firm must make them whole — and the data quality team has an uncomfortable Monday morning.',
    'At what stage in this process did the root cause failure occur? What system change would prevent it?'
)

make_key_concept(prs, ID, TITLE,
    'Corporate Action Types and Key Dates',
    [
        'MANDATORY: Cash Dividend, Bonus Issue, Stock Split, Mandatory Exchange Merger, Bankruptcy Distribution — happen automatically with no election required.',
        'VOLUNTARY: Rights Issue (exercise / sell nil-paid / lapse), Tender Offer (tender / retain), DRIP (take shares / take cash), Convertible Bond Conversion.',
        'Dividend Key Dates: Declaration Date → Ex-Dividend Date → Record Date → Payment Date. With T+1 settlement: Ex-date = Record date − 1 business day.',
        'Stock Split: 2-for-1 doubles share count, halves price per share. Total market cap unchanged. ALL position records must update simultaneously on ex-date.',
        'Back Office Process: identify eligible holders on record date → notify clients → receive and record elections → submit to registrar → process on effective date → update all position records.',
    ],
    '[TABLE: See Source Slide — Mandatory vs Voluntary corporate actions matrix with examples and operations impact]'
)

make_going_deeper(prs, ID, TITLE,
    'Mergers, Operational Risks & Regulatory Reporting',
    [
        'M&A Corporate Action: cash takeover (fixed cash per share) OR share-for-share exchange (exchange ratio) OR mixed consideration. Operations must manage election, verify entitlements, debit old position, credit new position on effective date.',
        'Old security is retired in the security master on effective date; new security is created (or existing acquirer record updated).',
        'Missed Election Risk: client loses voluntary entitlement — firm is typically liable to make the client whole (client money rules).',
        'Wrong Entitlement Calculation: dividend credited on wrong quantity — causes P&L discrepancy requiring manual correction and reconciliation investigation.',
        'FX Risk on Foreign Dividends: dividends paid in foreign currency must be converted at the correct rate on payment date — errors credited in wrong currency require reversal and reprocessing.',
        'Regulatory Reporting: all corporate action events must be correctly reflected in SEBI/MiFID II transaction and position reports.',
    ]
)

make_think_discuss(prs, ID, TITLE, [
    'A client misses a rights issue election deadline because they were travelling and did not check their email. The firm sent the notification to the correct address. Who is responsible — the firm or the client? What does your answer depend on?',
    'Stock splits are described as "purely mechanical" — the market cap doesn\'t change, only the share count. Yet operations teams flag them as high-risk events. Why?',
    'A large custody bank processes corporate actions for hundreds of thousands of client positions across 40 countries simultaneously. What does the operating model need to look like to manage this reliably?',
])

make_interview_ready(prs, ID, TITLE, [
    ('What is the difference between a mandatory and voluntary corporate action?',
     'Mandatory: happens automatically to all eligible holders — no election required (e.g., cash dividend, stock split, mandatory merger). Voluntary: holder must make an election (e.g., rights issue, tender offer, DRIP). Missing a voluntary deadline results in the default option applying.'),
    ('What are the four key dates in a dividend corporate action?',
     'Declaration Date (announced), Ex-Dividend Date (buy before this to receive dividend — T−1 from record date), Record Date (official holder register compiled), Payment Date (cash credited to accounts).'),
    ('What are the key operational risks in corporate actions processing?',
     'Missed election deadlines (client compensation liability), wrong entitlement calculations, FX conversion errors on foreign dividends, failure to update all position records simultaneously on ex-date, and incorrect regulatory reporting.'),
])

# ── MODULE COMPLETE SLIDE ──────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide, 0, 5.5, 13.33, 2.0, TEAL)
add_text_box(slide, 'MODULE 8 COMPLETE', 0.5, 0.5, 12.0, 0.8,
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text_box(slide, 'Securities Trading Organisation & RDM',
             0.5, 1.3, 12.0, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide,
    'STO.1 — STO Structure & Overview\n'
    'STO.2 — Front Office Operations\n'
    'STO.3 — Middle Office Operations\n'
    'STO.4 — Back Office & Settlement\n'
    'STO.5 — Reference Data Management\n'
    'STO.6 — Corporate Actions Processing',
    1.5, 2.5, 10.0, 3.0, size=15, color=LIGHT, align=PP_ALIGN.LEFT)
add_text_box(slide, 'The STO is the operational engine of capital markets. Every trade that is executed, confirmed, settled, and reported runs through this machine.',
             0.5, 5.65, 12.0, 1.2, size=13, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f'Saved: {OUTPUT_PATH}')
import os
print(f'File size: {os.path.getsize(OUTPUT_PATH):,} bytes')
print('Learning Deck DONE.')
